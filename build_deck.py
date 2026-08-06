#!/usr/bin/env python3
"""Codely · AI Summer Bootcamp Round 2 · Lesson 01 slide deck."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFilter

OUT = os.path.dirname(os.path.abspath(__file__))
W_IN, H_IN = 13.333, 7.5

# ---------- brand ----------
NAVY   = RGBColor(0x1E, 0x29, 0x3B)
SLATE  = RGBColor(0x64, 0x74, 0x8B)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
LIME   = RGBColor(0x84, 0xCC, 0x16)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)
BLUE   = RGBColor(0x38, 0x7C, 0xF0)
HEAD_F = "Lato"
BODY_F = "Lato"

# ---------- backgrounds ----------
def blob(img, cx, cy, r, color, alpha=255):
    layer = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(layer)
    d.ellipse([cx - r, cy - r, cx + r, cy + r], fill=color + (alpha,))
    layer = layer.filter(ImageFilter.GaussianBlur(r * 0.55))
    img.alpha_composite(layer)

def make_bg(path, peach=False):
    W, H = 2666, 1500
    img = Image.new("RGBA", (W, H), (255, 255, 255, 255))
    blob(img, int(W * 1.02), int(H * -0.08), 420, (220, 246, 145), 165)   # lime top-right
    blob(img, int(W * -0.02), int(H * 1.08), 395, (201, 169, 245), 150)   # purple bottom-left
    if peach:
        layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        d = ImageDraw.Draw(layer)
        r, cx, cy = 300, int(W * .805), int(H * .555)
        d.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(251, 214, 183, 255))
        layer = layer.filter(ImageFilter.GaussianBlur(2))
        img.alpha_composite(layer)
    img.convert("RGB").save(path, "PNG")

BG      = os.path.join(OUT, "_bg.png")
BG_SEC  = os.path.join(OUT, "_bg_section.png")
make_bg(BG)
make_bg(BG_SEC, peach=True)

# ---------- deck helpers ----------
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W_IN), Inches(H_IN)
BLANK = prs.slide_layouts[6]

def flat(shape):
    """Remove theme style + effects so LibreOffice renders no drop shadow."""
    shape.shadow.inherit = False
    sp = shape._element
    for st in sp.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}style'):
        sp.remove(st)
    return shape

def slide(section=False, logo=True):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(BG_SEC if section else BG, 0, 0, Inches(W_IN), Inches(H_IN))
    if logo:
        add_logo(s)
    return s

def add_logo(s):
    hexa = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(0.42), Inches(0.34), Inches(0.46), Inches(0.40))
    hexa.fill.solid(); hexa.fill.fore_color.rgb = GREEN
    hexa.line.fill.background(); flat(hexa)
    tf = _tf(hexa); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "</>"
    r.font.size, r.font.bold = Pt(9), True
    r.font.color.rgb, r.font.name = WHITE, BODY_F
    box(s, 0.96, 0.36, 1.6, 0.4, "CODELY", 13, GREEN, bold=True, space=True)

def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def box(s, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=None, space=False, line_spacing=None):
    sh = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _tf(sh); tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size, r.font.bold = Pt(size), bold
        r.font.color.rgb = color
        r.font.name = font or (HEAD_F if bold else BODY_F)
        if space:
            from pptx.oxml.ns import qn
            r.font._rPr.set("spc", "120")
    return sh

def title(s, text, size=40, y=0.85):
    return box(s, 0.85, y, 11.6, 1.5, text, size, NAVY, bold=True)

def card(s, x, y, w, h, border, head, body, head_size=17, body_size=12.5, fill=WHITE):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.09
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1.75)
    flat(sh)
    _tf(sh).text = ""
    box(s, x + 0.22, y + 0.26, w - 0.44, 0.5, head, head_size, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    box(s, x + 0.22, y + 0.95, w - 0.44, h - 1.2, body, body_size, SLATE,
        align=PP_ALIGN.CENTER, line_spacing=1.35)
    return sh

def pill(s, x, y, w, h, text, fill, tcolor=WHITE, size=13, radius=0.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); flat(sh)
    tf = _tf(sh); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(size), True, tcolor, HEAD_F
    return sh

def rect(s, x, y, w, h, text, border, size=14, fill=WHITE, tcolor=NAVY, bold=True, radius=0.12):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1.6)
    flat(sh)
    tf = _tf(sh); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.2
        r = p.add_run(); r.text = ln
        r.font.size, r.font.bold, r.font.color.rgb = Pt(size if i == 0 else size - 2), bold, tcolor
        r.font.name = HEAD_F if bold else BODY_F
    return sh

def line(s, x1, y1, x2, y2, color=LIGHT, width=1.6):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    return c

def dot(s, cx, cy, d, color):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); flat(sh)
    _tf(sh).text = ""
    return sh

def bullets(s, x, y, w, items, size=15, gap=0.62, dot_color=ORANGE):
    for i, t in enumerate(items):
        yy = y + i * gap
        dot(s, x + 0.09, yy + 0.16, 0.17, dot_color)
        box(s, x + 0.34, yy, w - 0.34, gap, t, size, NAVY, line_spacing=1.25)

def section(text, sub, tag="Build the future today!"):
    s = slide(section=True)
    box(s, 0.9, 2.42, 11.5, 1.5, text, 56, NAVY, bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.9, 3.98, 11.5, 0.6, sub, 20, SLATE, align=PP_ALIGN.CENTER)
    box(s, 0.9, 4.88, 11.5, 0.5, ">_  " + tag, 13, NAVY, bold=True, align=PP_ALIGN.CENTER)
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# =====================================================================
# 01 · Title
# =====================================================================
s = slide(logo=False)
hexa = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(6.09), Inches(1.62), Inches(1.15), Inches(1.0))
hexa.fill.solid(); hexa.fill.fore_color.rgb = GREEN
hexa.line.fill.background(); flat(hexa)
tf = _tf(hexa); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "</>"
r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(20), True, WHITE, BODY_F
box(s, 0.9, 2.95, 11.5, 1.3, "AI Summer Bootcamp", 62, NAVY, bold=True, align=PP_ALIGN.CENTER)
box(s, 0.9, 4.25, 11.5, 0.7, "Round 2  ·  Agentic AI + Мобайл апп бүтээх", 22, SLATE, align=PP_ALIGN.CENTER)
pill(s, 5.05, 5.15, 3.25, 0.5, "ХИЧЭЭЛ 1  ·  SETUP DAY", ORANGE)
notes(s, "0–10 мин. Танилцах. Round 2 эхэлж байна гэдгийг зарлах. Энергитэй!")

# =====================================================================
# 02 · Meet your instructor (placeholder)
# =====================================================================
s = slide()
title(s, "Meet Your Instructor")
ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(2.25), Inches(3.5), Inches(3.9))
ph.adjustments[0] = 0.05
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
ph.line.color.rgb = LIGHT; ph.line.width = Pt(1.5); flat(ph)
tf = _tf(ph); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "[ ЗУРАГ ]"
r.font.size, r.font.color.rgb, r.font.name = Pt(15), SLATE, BODY_F
box(s, 5.1, 2.45, 7.2, 0.7, "Б.Баясгалан", 30, ORANGE, bold=True)
bullets(s, 5.1, 3.35, 7.2, [
    "Co-founder @ Codely Academy",
    "МУИС-ийн ПХ инженер",
    "2015 оноос хойш 70+ төсөл",
    "MCS, Khan Bank, Andex IT — 10 жилийн туршлага",
], size=15, gap=0.62)
notes(s, "Хурдан. 2 минутаас хэтрүүлэхгүй. Photo placeholder-ыг зургаараа сольж тавь.")

# =====================================================================
# 03 · Mentor buddy
# =====================================================================
s = slide()
title(s, "Find Your Mission Partner")
box(s, 0.85, 1.95, 11.6, 0.5,
    "Round 2-т хос болж ажиллана: нэг нь ЖОЛООЧ (бичнэ), нөгөө нь НАВИГАТОР (хардаг). 15 мин тутам солино.",
    14, SLATE)
card(s, 0.95, 2.75, 3.5, 3.4, RGBColor(0xFB, 0xBF, 0x24), "Look Around!",
     "Хажууд суугаа хүүхэд бол\nтаны ХАМТРАГЧ.\nНэг баг гэсэн үг.")
card(s, 4.9, 2.75, 3.5, 3.4, PURPLE, "Mentor Buddy",
     "Round 1 төгссөн хүн\nшинэ найздаа тусална.\nТусалбал ✦ оноо!")
card(s, 8.85, 2.75, 3.5, 3.4, LIME, "Double Power",
     "Хоёулаа нэг компьютер,\nхоёулаа өөрийн утас.\nХамтдаа хурдан.")
notes(s, "Хосуудыг УРЬДЧИЛАН бэлдсэн жагсаалтаар зарла — тайзан дээр бодож болохгүй, цаг иднэ.\n"
         "Round 1 төгсөгчид 'Mentor buddy' гэсэн албан үүрэг өг. Discord-д нэгдүүл.")

# =====================================================================
# 04 · Section: Утсаа гарга
# =====================================================================
s = section("Утсаа гаргаарай", "Юу ч суулгаагүй байхад апп ажиллана")
notes(s, "10–25 мин. WOW #1. Энэ мөч өдрийн хамгийн чухал нь.")

# =====================================================================
# 05 · Expo Go + QR
# =====================================================================
s = slide()
title(s, "3 алхам, 5 минут")
for i, (h, b, c) in enumerate([
    ("1 · Суулга", "Play Store / App Store →\n\"Expo Go\"", PURPLE),
    ("2 · Нээ", "Expo Go →\nScan QR code", ORANGE),
    ("3 · Унш", "Дэлгэц дээрх\nQR-ыг унш", LIME),
]):
    card(s, 0.95 + i * 3.95, 2.35, 3.5, 2.35, c, h, b, head_size=18, body_size=13.5)
qr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.67), Inches(5.0), Inches(2.0), Inches(1.6))
qr.adjustments[0] = 0.06
qr.fill.solid(); qr.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
qr.line.color.rgb = NAVY; qr.line.width = Pt(2); flat(qr)
tf = _tf(qr); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "[ QR ЗУРАГ ]"
r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(14), True, NAVY, HEAD_F
box(s, 0.85, 6.72, 11.6, 0.5, "Anime Library  ·  багшийн бэлэн апп", 16, NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, "QR-ыг урьдчилан 2 өөр утсан дээр ТУРШСАН байх. Нэг Wi-Fi эсэхийг шалга.\n"
         "Ажиллахгүй бол: npx expo start --tunnel")

# =====================================================================
# 06 · Realtime WOW
# =====================================================================
s = slide()
box(s, 0.85, 1.4, 11.6, 1.2, "Одоо бүгд утсаа хараарай", 44, NAVY, bold=True, align=PP_ALIGN.CENTER)
box(s, 0.85, 2.65, 11.6, 0.6, "Би компьютер дээрээсээ шинэ кино нэмнэ.", 20, SLATE, align=PP_ALIGN.CENTER)
rect(s, 1.25, 3.75, 3.3, 1.5, "Багшийн вэб\nхуудас", PURPLE, size=17)
rect(s, 5.02, 3.75, 3.3, 1.5, "Firestore\n(database)", ORANGE, size=17)
rect(s, 8.79, 3.75, 3.3, 1.5, "Чиний утас", LIME, size=17)
box(s, 4.55, 4.28, 0.47, 0.5, "→", 26, SLATE, bold=True, align=PP_ALIGN.CENTER)
box(s, 8.32, 4.28, 0.47, 0.5, "→", 26, SLATE, bold=True, align=PP_ALIGN.CENTER)
pill(s, 4.3, 5.7, 4.75, 0.62, "2 СЕКУНД", ORANGE, size=18)
box(s, 0.85, 6.5, 11.6, 0.5, "Энэ бол Хичээл 4-т чиний бүтээх зүйл.", 17, NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Хамгийн хүчтэй мөч. Хүүхдүүд утсаа хараад байхад админ вэб дээрээ кино нэм.\n"
         "16 утсан дээр зэрэг гарч ирнэ. Урьдчилан 2 удаа турш!")

# =====================================================================
# 07 · Section: Юу сурах бол
# =====================================================================
s = section("Юу бүтээх бол?", "10 хичээлийн зам")

# =====================================================================
# 08 · 2 Season timeline
# =====================================================================
s = slide()
title(s, "The 10-Lesson Quest")
line(s, 1.15, 3.55, 12.15, 3.55, LIGHT, 3)
for cx, col, lbl, sub in [
    (2.35, PURPLE, "Season 1", "Хичээл 2–5\nБагштай хамт\nAnime Library"),
    (6.65, ORANGE, "Season 2", "Хичээл 6–9\nБие даан\nӨӨРИЙН апп"),
    (10.95, LIME, "Demo Day", "Хичээл 10\nЭцэг эх утсандаа\nтатаж туршина"),
]:
    dot(s, cx, 3.55, 0.32, col)
    box(s, cx - 1.9, 2.55, 3.8, 0.5, lbl, 22, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    box(s, cx - 1.9, 3.95, 3.8, 1.6, sub, 14.5, NAVY, align=PP_ALIGN.CENTER, line_spacing=1.35)
box(s, 0.85, 6.15, 11.6, 0.6,
    "Хичээл 1 (өнөөдөр) = хэрэгслээ бэлдэх.  Хичээл 2-оос бүтээж эхэлнэ.",
    15, SLATE, align=PP_ALIGN.CENTER)
notes(s, "Season 1 = багштай хамт нэг ижил бүтээгдэхүүн. Season 2 = өөрийн санаа.\n"
         "Season 1-ийн prompt-ууд чинь Season 2-д дахин хэрэглэгдэнэ гэдгийг ЗААВАЛ хэл.")

# =====================================================================
# 09 · Architecture
# =====================================================================
s = slide()
title(s, "Season 1: Anime Library", 36)
box(s, 0.85, 1.85, 11.6, 0.5, "Нэг бүтээгдэхүүн, 3 хэсэг — Netflix-тэй ижил бүтэц.", 15, SLATE)
rect(s, 1.1, 2.7, 4.4, 2.05,
     "WEB ADMIN\n\nкино нэмэх / засах / устгах\n(Google AI Studio)", PURPLE, size=16, radius=0.08)
rect(s, 7.85, 2.7, 4.4, 2.05,
     "MOBILE APP\n\nжагсаалт / дэлгэрэнгүй /\nхайлт / миний list  (Expo Go)", LIME, size=16, radius=0.08)
rect(s, 4.55, 5.35, 4.2, 1.35, "FIRESTORE\nнэг дата, 2 хэрэглэгч", ORANGE, size=16, radius=0.1)
line(s, 3.3, 4.75, 5.4, 5.35, SLATE, 2)
line(s, 10.05, 4.75, 7.95, 5.35, SLATE, 2)
notes(s, "Гол ойлголт: контент нэмдэг хүн (admin) vs контент үздэг хүн (user). 2 өөр интерфэйс, НЭГ дата.\n"
         "Web admin-ыг AI Studio-гаар хийнэ → Round 1 төгсөгчид энд Mentor болно.")

# =====================================================================
# 10 · Season 1 lessons
# =====================================================================
s = slide()
title(s, "Season 1 — Хичээл 2–5")
items = [
    ("Хичээл 2", "Кинонуудын жагсаалт + дэлгэрэнгүй дэлгэц", PURPLE),
    ("Хичээл 3", "Админ вэб хуудас — кино нэмдэг, засдаг", ORANGE),
    ("Хичээл 4", "Апп ↔ Firestore: realtime шидэт мөч", LIME),
    ("Хичээл 5", "Камер · хайлт · Миний жагсаалт · share", BLUE),
]
for i, (h, b, c) in enumerate(items):
    y = 2.15 + i * 1.15
    pill(s, 0.95, y + 0.12, 1.85, 0.62, h, c, size=15)
    box(s, 3.15, y + 0.2, 9.2, 0.6, b, 18, NAVY, bold=True)
    if i < 3:
        line(s, 0.95, y + 1.02, 12.35, y + 1.02, LIGHT, 1.2)
pill(s, 0.95, 6.75, 4.1, 0.5, "→  БЭЛЭН БҮТЭЭГДЭХҮҮН", GREEN, size=13)
notes(s, "Х5-ын төгсгөлд бүх хүүхэд бүрэн ажилладаг бүтээгдэхүүнтэй болно: вэб + апп + database.")

# =====================================================================
# 11 · Season 2
# =====================================================================
s = slide()
title(s, "Season 2 — ӨӨРИЙН апп")
box(s, 0.85, 1.85, 11.6, 0.5,
    "Гэр бүл, найзаасаа бодит асуудал \"агнаж\" ирээд түүнийг шийддэг апп бүтээнэ.", 15, SLATE)
for i, (h, b, c) in enumerate([
    ("Хичээл 6", "Problem Hunt\nӨөрийн санаа + төлөвлөгөө", PURPLE),
    ("Хичээл 7", "Build sprint #1\nГол дэлгэц + database", ORANGE),
    ("Хичээл 8", "Build sprint #2\nMaps · сануулга · polish", LIME),
    ("Хичээл 9", "Тест + Pitch\nБодит хэрэглэгчээр туршуулах", BLUE),
]):
    card(s, 0.95 + i * 2.92, 2.7, 2.6, 2.5, c, h, b, head_size=16, body_size=12.5)
box(s, 0.85, 5.6, 11.6, 1.1,
    "Х2–5-д хэрэглэсэн бүх prompt-оо цуглуулж явна = чиний PLAYBOOK.\n"
    "Шинэ аппыг 0-оос биш, загвартай эхэлнэ.", 16, NAVY, bold=True,
    align=PP_ALIGN.CENTER, line_spacing=1.35)
notes(s, "Playbook гэдгийг тодорхой тайлбарла — Season 1 хаягдахгүй, дахин хэрэглэгдэнэ.")

# =====================================================================
# 12 · Demo Day
# =====================================================================
s = slide()
title(s, "Хичээл 10: Demo Day")
card(s, 0.95, 2.35, 3.5, 3.5, PURPLE, "The Final Build",
     "Аппаа бэлэн болгож,\nicon, нэр, splash\nхийж дуусгана.")
card(s, 4.9, 2.35, 3.5, 3.5, ORANGE, "Эцэг эхийн утсанд",
     "Эцэг эх чинь\nExpo Go-оор чиний аппыг\nӨӨРИЙН утсандаа татаж\nтуршина.")
card(s, 8.85, 2.35, 3.5, 3.5, LIME, "Glory & Prizes",
     "Best Problem Solver\nBest Design\nBest Pitch\nBest Mentor Buddy")
box(s, 0.85, 6.25, 11.6, 0.5, "8 сарын 22  ·  Баасан", 16, SLATE, align=PP_ALIGN.CENTER)
notes(s, "Demo Day-г ЭХНИЙ өдрөөс зарлах нь motivation-ы гол хөшүүрэг. Огноог тодорхой хэл.")

# =====================================================================
# 13 · Round 1 vs Round 2
# =====================================================================
s = slide()
title(s, "Round 1 → Round 2")
hdr = [("", 0.95, 2.6), ("Round 1", 4.15, 3.7), ("Round 2", 8.35, 4.0)]
for t, x, w in hdr:
    if t: box(s, x, 2.1, w, 0.5, t, 19, ORANGE, bold=True)
rows = [
    ("Гарц", "Вэб сайт, тоглоом", "Вэб + утасны апп + database"),
    ("AI", "Чат (AI Studio)", "Agent — файл засаж, ажиллуулна"),
    ("Ур чадвар", "Prompt бичих", "Agent-ыг удирдах"),
    ("Санаа", "\"Ямар тоглоом?\"", "\"Хэний ямар асуудал?\""),
    ("Demo Day", "Дэлгэц дээр", "Эцэг эхийн утсан дээр"),
]
for i, (a, b, c) in enumerate(rows):
    y = 2.85 + i * 0.78
    box(s, 0.95, y, 3.0, 0.6, a, 15, SLATE, bold=True)
    box(s, 4.15, y, 3.9, 0.6, b, 15, SLATE)
    box(s, 8.35, y, 4.0, 0.6, c, 15, NAVY, bold=True)
    line(s, 0.95, y + 0.66, 12.35, y + 0.66, LIGHT, 1)
notes(s, "Шинэ сурагчид: 'Round 1 хийгээгүй байсан ч гүйцнэ' гэдгийг тодорхой хэл.\n"
         "Round 1 төгсөгчид: 'AI Studio-гийн мэдлэг чинь Х3-т хэрэгтэй болно' — бахархал өг.")

# =====================================================================
# 14 · Section: Апп хаанаас ирдэг вэ
# =====================================================================
s = section("Апп хаанаас ирдэг вэ?", "Апп хөгжүүлэлтийн ертөнц")
notes(s, "40–55 мин. Хамгийн 'лекц' хэсэг — 15 минутаас ХЭТРҮҮЛЭХГҮЙ. Асуулт асууж оролцуул.")

# =====================================================================
# 15 · 3 kinds of apps
# =====================================================================
s = slide()
title(s, "3 төрлийн апп")
card(s, 0.95, 2.2, 3.5, 3.5, PURPLE, "Web app",
     "Браузераар нээдэг.\nСуулгах шаардлагагүй.\n\nRound 1-д хийсэн сайт")
card(s, 4.9, 2.2, 3.5, 3.5, RGBColor(0xEF, 0x44, 0x44), "Native app",
     "Утсанд суудаг.\niPhone-д Swift,\nAndroid-д Kotlin —\n2 өөр код бичнэ.")
card(s, 8.85, 2.2, 3.5, 3.5, LIME, "Cross-platform",
     "1 код → 2 платформ.\n\nManай зам:\nReact Native + Expo")
box(s, 0.85, 6.05, 11.6, 0.7, "Асуулт: утсан дээрээ хамгийн их ашигладаг апп чинь аль төрөл вэ?",
    16, NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Асуултаар оролцуул. Instagram, TikTok = native (эсвэл cross-platform). Хүүхдүүд таамаглана.")

# =====================================================================
# 16 · Why Expo
# =====================================================================
s = slide()
title(s, "Яагаад Expo?")
bullets(s, 1.05, 2.3, 11.3, [
    "1 удаа бичээд iPhone, Android хоёуланд нь ажиллана",
    "Expo Go — App Store-т тавихгүйгээр утсан дээрээ шууд туршина",
    "Код өөрчлөгдөх бүрт утас чинь 2–3 секундэд шинэчлэгдэнэ",
    "Камер, зураг, GPS, сануулга — утасны бүх боломж бэлэн байдаг",
], size=18, gap=0.85)
box(s, 1.05, 5.95, 11.3, 0.8,
    "Тиймээс өнөөдөр Expo Go-г утсандаа суулгана.", 19, NAVY, bold=True)
notes(s, "Round 1-д Vercel дээр deploy хийсэн шиг — Expo Go бол утасны 'deploy'.")

# =====================================================================
# 17 · Roles
# =====================================================================
s = slide()
title(s, "Апп бүтээх багт хэн байдаг вэ?")
roles = [
    ("Product owner", "Юу бүтээх,\nхэнд зориулах", PURPLE),
    ("Designer", "Дэлгэц ямар\nхарагдахыг зурна", ORANGE),
    ("Developer", "Код бичнэ", RGBColor(0xEF, 0x44, 0x44)),
    ("Tester", "Эвдэрч байна уу\nгэж шалгана", LIME),
]
for i, (h, b, c) in enumerate(roles):
    card(s, 0.95 + i * 2.92, 2.2, 2.6, 2.3, c, h, b, head_size=16, body_size=13)
box(s, 0.95, 5.0, 11.4, 1.6,
    "ЧИ = Product owner + Designer + Tester\nЧИНИЙ AGENT = Developer",
    26, NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.4)
notes(s, "Энэ бол Round 2-ийн гол framing. Хүүхэд код бичихгүй — удирдана, шалгана.")

# =====================================================================
# 18 · Section: Agentic AI
# =====================================================================
s = section("Agentic AI гэж юу вэ?", "Чатнаас юугаараа өөр вэ?")
notes(s, "60–75 мин. Live demo бэлэн байг.")

# =====================================================================
# 19 · Chat vs Agent
# =====================================================================
s = slide()
title(s, "Чат AI  vs  Agent")
box(s, 4.5, 2.15, 3.8, 0.5, "Чат AI (Round 1)", 19, SLATE, bold=True, align=PP_ALIGN.CENTER)
box(s, 8.5, 2.15, 3.9, 0.5, "Agent (Round 2)", 19, ORANGE, bold=True, align=PP_ALIGN.CENTER)
for i, (a, b, c) in enumerate([
    ("Юу хийдэг", "Хариу бичиж өгнө", "Файлыг өөрөө засна"),
    ("Терминал", "Ашиглаж чадахгүй", "Команд ажиллуулна"),
    ("Шалгалт", "Чи хуулж тавина", "Өөрөө ажиллуулж шалгана"),
    ("Алдаа", "Чи олно", "Алдаагаа өөрөө уншиж заслаа"),
]):
    y = 2.95 + i * 0.85
    box(s, 0.95, y, 3.3, 0.6, a, 16, NAVY, bold=True)
    box(s, 4.5, y, 3.8, 0.6, b, 15, SLATE, align=PP_ALIGN.CENTER)
    box(s, 8.5, y, 3.9, 0.6, c, 15, NAVY, bold=True, align=PP_ALIGN.CENTER)
    line(s, 0.95, y + 0.68, 12.35, y + 0.68, LIGHT, 1)
notes(s, "Live demo: 1 prompt өгөөд agent файл үүсгэж, команд ажиллуулж байгааг проектор дээр харуул.")

# =====================================================================
# 20 · The loop
# =====================================================================
s = slide()
title(s, "Агентын 4 алхамын гогцоо")
box(s, 0.85, 1.9, 11.6, 0.5, "Хичээл бүрт, feature бүрт — үүнийг л дахин дахин хийнэ.", 15, SLATE)
steps = [
    ("1 · PROMPT", "Юу хүсэж байгаагаа\ntодорхой хэл", PURPLE),
    ("2 · BUILD", "Agent код бичнэ,\nутас 2 сек-д шинэчлэгдэнэ", ORANGE),
    ("3 · VERIFY", "УТСАН ДЭЭРЭЭ шалга\nхүссэн шиг болов уу?", LIME),
    ("4 · RE-PROMPT", "Болоогүй бол \"би юу\nхарж байна\"-г хэл", BLUE),
]
for i, (h, b, c) in enumerate(steps):
    card(s, 0.95 + i * 2.92, 2.65, 2.48, 2.4, c, h, b, head_size=15.5, body_size=12.5)
    if i < 3:
        box(s, 3.52 + i * 2.92, 3.6, 0.4, 0.5, "→", 22, SLATE, bold=True, align=PP_ALIGN.CENTER)
box(s, 0.95, 5.45, 11.4, 1.3,
    "Agent чиний утсыг МЭДРЭХ ЧАДВАРГҮЙ.\nТовч жижиг үү, чирэхэд гацаж байна уу — үүнийг зөвхөн ЧИ мэднэ.",
    18, NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.4)
notes(s, "Tester гэдэг үүрэг яг үүнээс гардаг. Хүүхдэд 'чи л мэднэ' гэдгийг бахархал болгож өг.")

# =====================================================================
# 21 · Section: Setup
# =====================================================================
s = section("Setup marathon", "Хосоороо · 45 минут · checklist-ээр")
notes(s, "75–120 мин. Хамгийн гацдаг хэсэг. Дуусгасан хос → Mentor болж бусдад тусална (✦ оноо).")

# =====================================================================
# 22 · Setup steps
# =====================================================================
s = slide()
title(s, "8 алхам")
left = [("A", "Утсандаа Expo Go"), ("B", "Node.js"), ("C", "Agent (Codex / Antigravity)"), ("D", "Expo account")]
right = [("E", "Expo Skills"), ("F", "Expo MCP"), ("G", "Шалгах prompt"), ("H", "Эхний өөрийн апп")]
for col, items in enumerate([left, right]):
    for i, (k, t) in enumerate(items):
        y = 2.25 + i * 1.0
        x = 0.95 + col * 6.1
        pill(s, x, y + 0.05, 0.62, 0.62, k, ORANGE if col == 0 else PURPLE, size=17)
        box(s, x + 0.85, y + 0.13, 4.9, 0.6, t, 17, NAVY, bold=True)
box(s, 0.95, 6.3, 11.4, 0.6,
    "Гацвал: алдааны текстийг БҮТНЭЭР хуулж agent-даа өг. Тэр өөрөө засна.",
    16, NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Checklist хуудсыг хос бүрт тарааж бай. Самбарт 'бэлэн болсон хосууд' жагсаалт бич — өрсөлдөөн үүснэ.\n"
         "Гацсан сурагчийг 10 минутаас илүү барихгүй → нөөц машин руу шилжүүл.")

# =====================================================================
# 23 · Skills + MCP
# =====================================================================
s = slide()
title(s, "Агентдаа мэдлэг нэмнэ")
box(s, 0.85, 1.9, 11.6, 0.5, "Agent анхнаасаа Expo-г мэдэхгүй. Бид түүнийг ухаалаг болгоно.", 15, SLATE)
rect(s, 1.1, 2.75, 4.9, 2.5,
     "EXPO SKILLS\n\nExpo-г яаж зөв бүтээхийг\nагентад заана", PURPLE, size=17, radius=0.08)
rect(s, 7.35, 2.75, 4.9, 2.5,
     "EXPO MCP\n\nАгент албан ёсны документыг\nӨӨРӨӨ уншина", LIME, size=17, radius=0.08)
box(s, 0.95, 5.6, 11.4, 1.2,
    "Round 2-ийн гол ур чадвар: агентыг зүгээр хэрэглэхгүй — түүнийг ХҮЧТЭЙ БОЛГОНО.",
    18, NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.35)
notes(s, "Энэ бол Round 2-ийн үнэ цэнийн гол мөч: prompt бичигчээс → agent удирдагч болох.")

# =====================================================================
# 24 · AI Smart & Safe
# =====================================================================
s = slide()
title(s, "AI Smart & Safe: \"Зөвшөөрөх үү?\"")
box(s, 0.85, 1.9, 11.6, 0.5, "Agent ажиллахдаа зөвшөөрөл асууна. Хаа сайгүй Yes дарж болохгүй.", 15, SLATE)
rows = [
    ("Төслийн файлыг засах уу?", "ТИЙМ", GREEN),
    ("npm install ... ажиллуулах уу?", "ТИЙМ", GREEN),
    ("Файл УСТГАХ уу?", "УНШААД", RGBColor(0xFB, 0xBF, 0x24)),
    ("Бүх зөвшөөрлийг автоматаар өгөх үү?", "ҮГҮЙ", RGBColor(0xEF, 0x44, 0x44)),
    ("Төслөөс ГАДУУРХ файлд гар хүрэх", "ҮГҮЙ", RGBColor(0xEF, 0x44, 0x44)),
]
for i, (q, a, c) in enumerate(rows):
    y = 2.7 + i * 0.72
    box(s, 0.95, y + 0.05, 8.3, 0.6, q, 16, NAVY)
    pill(s, 9.6, y, 2.4, 0.56, a, c, size=14)
box(s, 0.95, 6.5, 11.4, 0.6, "Дүрэм: ойлгоогүй зүйл дээр Yes дарахгүй. Асуу.",
    18, NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, "135–145 мин. Богино, тодорхой. Жишээ түүх ярь: auto-approve дараад файлаа алдсан хүн.")

# =====================================================================
# 25 · Homework
# =====================================================================
s = slide()
title(s, "Гэрийн даалгавар")
bullets(s, 1.05, 2.35, 11.3, [
    "Утсан дээрх аппаа гэр бүлдээ үзүүл — \"энэ намайг хийсэн\"",
    "Дуртай 5 кино / аниме / манга-гийн нэрийг Discord-д бич  (Хичээл 2-д хэрэглэнэ)",
    "Дуртай 3 аппыг бод: тэд ямар асуудлыг шийдэж байна вэ?",
], size=18, gap=0.95)
pill(s, 0.95, 5.5, 6.7, 0.72, "ДАРААГИЙН ХИЧЭЭЛД УТСАА ЗААВАЛ АВЧИР", ORANGE, size=16)
box(s, 0.95, 6.45, 11.4, 0.5, "Хичээл 2  ·  8 сарын 4, Мягмар  ·  09:00", 16, SLATE)
notes(s, "145–150 мин. Утас, цэнэглэгч — эцэг эхэд бас Discord-оор сануул.")

# =====================================================================
# 26 · Closing
# =====================================================================
s = section("Асуулт хариулт?", "Ready to build your first app?")

# ---------- save ----------
pptx_path = os.path.join(OUT, "Codely-Round2-Lesson01-Roadmap.pptx")
prs.save(pptx_path)
print("saved:", pptx_path, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
